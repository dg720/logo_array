from pptx import Presentation
from pptx.util import Inches
from PIL import Image
import os

# Configuration
LOGO_FOLDER = "C:/Users/ASUS/Documents/logo_array/src/logos"  # C:/Users/dhruv\Documents/logo_array/src/logos"  # Folder where logos are stored
OUTPUT_FILE = "logos_presentation.pptx"
# LOGO_HEIGHT = 15  # Standardized height in pixels
LOGO_POSITIONS = (3, 10)  # (columns, rows) grid structure # ONLY INPUT
USER_WIDTH, USER_HEIGHT = (4, 9)  # user input dim in inches # ONLY INPUT
SLIDE_WIDTH, SLIDE_HEIGHT = (
    Inches(USER_WIDTH),
    Inches(USER_HEIGHT),
)  # Standard slide size
HORIZONTAL_PADDING = Inches(2)  # Extra spacing between columns
WHITE_THRESHOLD = 230  # Adjust threshold for detecting white backgrounds (0-255)
# MAX_WIDTH = 120

num_cols, num_rows = LOGO_POSITIONS[0], LOGO_POSITIONS[1]
LOGO_HEIGHT = int(USER_HEIGHT * 96 / num_rows / 2)
slide_pixel_width = USER_WIDTH * 96
MAX_WIDTH = int(slide_pixel_width / num_cols)
# HORIZONTAL_PADDING = Inches(slide_pixel_width / num_cols * 1 / 3 / 96)

# Load and process logos
logos = [f for f in os.listdir(LOGO_FOLDER) if f.endswith((".png", ".jpg", ".jpeg"))]
processed_logos = []


def remove_white_background(image):
    """Converts white background pixels to transparent."""
    image = image.convert("RGBA")  # Ensure image is in RGBA mode
    data = image.getdata()

    new_data = []
    for item in data:
        # Check if pixel is close to white
        if (
            item[0] > WHITE_THRESHOLD
            and item[1] > WHITE_THRESHOLD
            and item[2] > WHITE_THRESHOLD
        ):
            new_data.append((255, 255, 255, 0))  # Make transparent
        else:
            new_data.append(item)  # Keep original pixel

    image.putdata(new_data)
    return image


def auto_crop(image):
    """Auto-crops an image by removing fully transparent borders."""
    bbox = image.getbbox()
    if bbox:
        return image.crop(bbox)
    return image  # If no cropping needed, return original image


# Resize and remove background
for logo in logos:
    logo_path = os.path.join(LOGO_FOLDER, logo)
    img = Image.open(logo_path)

    # Step 1: Remove white background
    img = remove_white_background(img)

    # Step 2: Auto-crop to remove excess transparent pixels
    img = auto_crop(img)

    # Step 3: Resize while maintaining aspect ratio
    aspect_ratio = img.width / img.height
    new_width = int(LOGO_HEIGHT * aspect_ratio)
    new_height = LOGO_HEIGHT
    if new_width > MAX_WIDTH:
        new_width = MAX_WIDTH
        new_height = int(new_width / aspect_ratio)
    img = img.resize((new_width, new_height))

    # Overwrite the original image
    img.save(logo_path, format="PNG")  # Ensure PNG format for transparency

    processed_logos.append((logo_path, new_width, new_height))

## **Step 4: Compute the widest image per column**
# column_widths = [0] * LOGO_POSITIONS[0]  # Track max width in each column
# for idx, (_, width, _) in enumerate(processed_logos):
#    col = idx % LOGO_POSITIONS[0]  # Determine column index
#    column_widths[col] = max(column_widths[col], width)

# **Step 5: Compute column centers**
col_spacing = SLIDE_WIDTH / (LOGO_POSITIONS[0] - 1)
column_centers = []
current_x = Inches(0.1)  # Start slightly from the left margin
for i in range(num_cols):
    column_center = current_x + col_spacing / 2  # Center within column
    column_centers.append(column_center)
    current_x += col_spacing  # Move to next column start

# **Step 6: Create PowerPoint presentation**
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank slide

# **Step 7: Calculate vertical spacing**
row_spacing = SLIDE_HEIGHT / (LOGO_POSITIONS[1] - 1)  # count number of spaces not rows

# **Step 8: Place logos in a grid with center alignment**
for idx, (logo, width, height) in enumerate(processed_logos):
    col = idx % LOGO_POSITIONS[0]  # Determine column index
    row = idx // LOGO_POSITIONS[0]  # Determine row index

    if row >= LOGO_POSITIONS[1]:  # Stop if exceeding grid limit
        break

    # Convert pixel dimensions to inches (PowerPoint uses inches, assuming 96 DPI)
    width_inches = Inches(width / 96)
    height_inches = Inches(height / 96)

    # Center x position within its column
    x_pos = column_centers[col] - (width_inches / 2)  # Shift left by half width

    # Calculate y position (center in row)
    y_pos = (row + 1) * row_spacing - height_inches / 2

    # Insert logo as a picture object
    slide.shapes.add_picture(
        logo, x_pos, y_pos, width=width_inches, height=height_inches
    )

# Save PowerPoint
prs.save(OUTPUT_FILE)
print(f"PowerPoint saved as {OUTPUT_FILE}")
